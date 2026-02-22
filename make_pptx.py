"""Generate ChemEAGLE pipeline + failure analysis PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x0F, 0x17, 0x2A)   # near-black navy
MID_BG    = RGBColor(0x16, 0x24, 0x3E)   # panel background
ACCENT1   = RGBColor(0x00, 0xB4, 0xD8)   # cyan
ACCENT2   = RGBColor(0x48, 0xCA, 0xE4)   # light cyan
WARN      = RGBColor(0xFF, 0xA5, 0x00)   # amber / warning
DANGER    = RGBColor(0xFF, 0x4D, 0x6D)   # red
OK        = RGBColor(0x38, 0xD9, 0xA9)   # teal/green
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY     = RGBColor(0xB0, 0xC4, 0xDE)
MGRAY     = RGBColor(0x4A, 0x5E, 0x78)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helpers ─────────────────────────────────────────────────────────────────

def fill_bg(slide, color=DARK_BG):
    """Solid-fill slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color=MID_BG, line_color=None, line_w_pt=0):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.ROUNDED_RECTANGLE if False else 1,
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=16, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_para(tf, text, font_size=13, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p


def slide_header(slide, title, subtitle=None):
    """Top accent bar + title."""
    # accent bar
    bar = add_rect(slide, 0, 0, 13.33, 0.08, fill_color=ACCENT1)
    # title
    add_text(slide, title, 0.4, 0.12, 12.0, 0.65,
             font_size=30, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.72, 12.0, 0.38,
                 font_size=15, color=LGRAY, italic=True)
    # bottom bar
    add_rect(slide, 0, 7.42, 13.33, 0.08, fill_color=ACCENT1)
    add_text(slide, "ChemEAGLE  |  Pipeline & Failure Analysis",
             0.4, 7.42, 10.0, 0.08, font_size=8, color=DARK_BG, align=PP_ALIGN.LEFT)


def pill(slide, text, l, t, w, h, bg=ACCENT1, fg=DARK_BG, font_size=11, bold=True):
    r = add_rect(slide, l, t, w, h, fill_color=bg, line_color=None)
    add_text(slide, text, l+0.04, t+0.03, w-0.08, h-0.06,
             font_size=font_size, bold=bold, color=fg, align=PP_ALIGN.CENTER)
    return r


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
fill_bg(slide)

# large cyan accent top stripe
add_rect(slide, 0, 0, 13.33, 0.12, fill_color=ACCENT1)
# bottom stripe
add_rect(slide, 0, 7.38, 13.33, 0.12, fill_color=ACCENT1)

# decorative right panel
add_rect(slide, 9.8, 0.12, 3.53, 7.26, fill_color=MID_BG)

add_text(slide, "ChemEAGLE", 0.6, 1.5, 9.0, 1.2,
         font_size=60, bold=True, color=WHITE)
add_text(slide, "Multi-Agent Chemical Reaction Extraction", 0.6, 2.7, 9.0, 0.6,
         font_size=24, bold=False, color=ACCENT2)
add_text(slide, "Pipeline Architecture  ·  Agent Roles  ·  Failure Analysis",
         0.6, 3.35, 8.5, 0.5, font_size=16, color=LGRAY, italic=True)

# right panel content
add_text(slide, "10 Failure\nModes\nIdentified", 10.0, 1.8, 3.0, 1.6,
         font_size=28, bold=True, color=WARN, align=PP_ALIGN.CENTER)
add_text(slide, "5 Pipeline\nStages", 10.0, 3.6, 3.0, 1.0,
         font_size=28, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)
add_text(slide, "4 Execution\nPaths", 10.0, 4.8, 3.0, 1.0,
         font_size=28, bold=True, color=OK, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — System Overview
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
fill_bg(slide)
slide_header(slide, "System Overview",
             "What ChemEAGLE does and how it fits into a research workflow")

# 3-column tiles
tiles = [
    ("INPUT",   ACCENT1, "PDF files or\nindividual\nreaction images",
     "Florence-2 OD detects\nfigures & tables;\ncrops exported as PNG"),
    ("PROCESS", WARN,    "5-stage multi-agent\npipeline extracts\nchemical reactions",
     "LLM planner → specialist\ntools (RxnIM, MolNexTR,\nOCR) → LLM synthesizer"),
    ("OUTPUT",  OK,      "Structured JSON:\nreactants, products,\nconditions, SMILES",
     "Machine-readable reaction\ndatabase ready for\ndownstream analysis"),
]
for i, (label, color, top, bot) in enumerate(tiles):
    lx = 0.4 + i * 4.3
    add_rect(slide, lx, 1.25, 4.0, 5.6, fill_color=MID_BG,
             line_color=color, line_w_pt=2)
    add_text(slide, label, lx+0.1, 1.3, 3.8, 0.45,
             font_size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    # divider
    add_rect(slide, lx+0.15, 1.8, 3.7, 0.03, fill_color=color)
    add_text(slide, top, lx+0.15, 1.9, 3.7, 1.5,
             font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, bot, lx+0.15, 3.5, 3.7, 2.9,
             font_size=12, color=LGRAY, align=PP_ALIGN.CENTER)

# arrow connectors
for ax in [4.42, 8.72]:
    add_text(slide, "▶", ax, 3.6, 0.5, 0.5,
             font_size=28, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Agents & Models
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
fill_bg(slide)
slide_header(slide, "Agents & Models", "Every component involved in the extraction pipeline")

rows = [
    # (Role, Model/Tool, Type, Purpose)
    ("PDF Extractor",       "VisualHeist / Florence-2",       "Vision OD",       "Detects & crops all figures from PDF pages"),
    ("Planner",             "GPT-4o-mini / Qwen3-VL",         "Multimodal LLM",  "Classifies image type; selects which agents to run"),
    ("Plan Observer",       "GPT-4o-mini / Qwen3-VL",         "Multimodal LLM",  "Second-opinion sanity check on the plan"),
    ("Reaction Template",   "RxnIM  (rxn.ckpt)",              "Specialist CV",   "Detects reaction arrows; maps reactant/product regions"),
    ("Molecule Detector",   "ChemIEToolkit MolDetector",      "Specialist CV",   "Localises individual molecules within a scheme"),
    ("SMILES Generator",    "MolNexTR  (molnextr.ckpt)",      "Graph Neural Net","Image→graph→SMILES for each molecule crop"),
    ("Table Parser",        "TableParser + TesseractOCR",     "OCR / layout",    "Extracts R-group substitution tables from images"),
    ("Text NLP",            "ChemRxnExtractor + ChemNER",     "BioBERT NLP",     "Extracts reactions & named entities from caption text"),
    ("Condition Classifier","GPT-4o-mini / Qwen3-VL",         "LLM",             "Assigns solvent / reagent / catalyst roles to conditions"),
    ("Action Observer",     "GPT-4o-mini / Qwen3-VL",         "Multimodal LLM",  "Validates tool outputs; signals redo if results are bad"),
    ("Data Structurer",     "GPT-4o-mini / Qwen3-VL",         "LLM",             "Assembles final JSON; converts abbreviations to SMILES"),
]

col_w = [2.3, 2.5, 1.6, 5.5]
col_x = [0.3, 2.65, 5.2, 6.85]
headers = ["Role", "Model / Tool", "Type", "Purpose"]
header_y = 1.18
row_h = 0.44
row_y0 = 1.62

# header row
for ci, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    add_rect(slide, cx, header_y, cw-0.06, 0.38, fill_color=ACCENT1)
    add_text(slide, hdr, cx+0.05, header_y+0.04, cw-0.12, 0.3,
             font_size=12, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

for ri, (role, model, typ, purpose) in enumerate(rows):
    ry = row_y0 + ri * row_h
    bg = RGBColor(0x1A, 0x2C, 0x48) if ri % 2 == 0 else MID_BG
    add_rect(slide, col_x[0], ry, 12.57, row_h-0.04, fill_color=bg)
    vals = [role, model, typ, purpose]
    colors = [ACCENT2, WHITE, LGRAY, LGRAY]
    bolds  = [True, True, False, False]
    for ci, (val, cx, cw, col, bld) in enumerate(zip(vals, col_x, col_w, colors, bolds)):
        add_text(slide, val, cx+0.07, ry+0.05, cw-0.12, row_h-0.1,
                 font_size=10, bold=bld, color=col, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Pipeline Stage 1-2
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
fill_bg(slide)
slide_header(slide, "Pipeline: Stages 1 – 2",
             "Input ingestion → Planning → Plan validation")

stage_data = [
    ("STAGE 0\nInput", ACCENT1,
     ["PDF → VisualHeist (Florence-2 OD)", "Crops ALL figures to PNG",
      "Direct image input also accepted"]),
    ("STAGE 1\nPlanner", WARN,
     ["Multimodal LLM sees the image", "Classifies: reaction / R-group table / text",
      "Selects which tool agents to invoke", "Output: ordered agent_list"]),
    ("STAGE 2\nPlan Observer", ACCENT2,
     ["Second LLM reviews the plan", "Can revise or confirm agent selection",
      "Uses same model → shared blind spots", "Disabled by default in OS mode"]),
]

xs = [0.35, 4.55, 8.75]
for i, (title, color, bullets) in enumerate(stage_data):
    lx = xs[i]
    add_rect(slide, lx, 1.2, 4.0, 5.65, fill_color=MID_BG,
             line_color=color, line_w_pt=2)
    add_text(slide, title, lx+0.12, 1.26, 3.76, 0.75,
             font_size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_rect(slide, lx+0.12, 2.02, 3.76, 0.03, fill_color=color)
    for bi, b in enumerate(bullets):
        add_text(slide, f"• {b}", lx+0.18, 2.12 + bi*0.72, 3.65, 0.68,
                 font_size=12, color=WHITE if color != ACCENT2 else DARK_BG)

    if i < 2:
        add_text(slide, "▶", lx+4.05, 3.85, 0.4, 0.4,
                 font_size=20, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)

# warning note at bottom
add_rect(slide, 0.35, 6.92, 12.58, 0.42, fill_color=RGBColor(0x3D,0x1A,0x00),
         line_color=WARN, line_w_pt=1)
add_text(slide,
         "⚠  Key risk: both Planner and Plan Observer use the same LLM — "
         "a misclassification in Stage 1 is rarely caught in Stage 2.",
         0.5, 6.94, 12.3, 0.38, font_size=11, color=WARN)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Pipeline Stage 3 (4 paths)
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
fill_bg(slide)
slide_header(slide, "Pipeline: Stage 3 — Tool Executor",
             "Four specialist paths; selected by the Planner's agent_list")

path_data = [
    ("PATH A", ACCENT1, "Reaction Template",
     ["RxnIM detects arrow + regions", "MolDetector localises molecules",
      "MolNexTR → SMILES per crop", "LLM classifies condition roles"]),
    ("PATH B", OK, "Structure-based R-groups",
     ["PATH A first to get scaffold", "Each variant panel → MolNexTR",
      "normalize_product_variant_output()", "Scaffold + R substituents merged"]),
    ("PATH C", WARN, "Table-based R-groups",
     ["PATH A for scaffold SMILES", "TableParser detects table region",
      "TesseractOCR reads each cell", "LLM maps cells to R positions"]),
    ("PATH D", LGRAY, "Text Extraction",
     ["TesseractOCR on image caption", "ChemRxnExtractor (BioBERT NLP)",
      "ChemNER identifies entities", "Structured text reaction record"]),
]

xs = [0.3, 3.6, 6.9, 10.15]
ws = 3.12
for i, (path, color, subtitle, bullets) in enumerate(path_data):
    lx = xs[i]
    add_rect(slide, lx, 1.18, ws, 6.1, fill_color=MID_BG,
             line_color=color, line_w_pt=2)
    add_text(slide, path, lx+0.08, 1.22, ws-0.16, 0.38,
             font_size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, subtitle, lx+0.08, 1.6, ws-0.16, 0.38,
             font_size=10, bold=False, color=LGRAY, align=PP_ALIGN.CENTER, italic=True)
    add_rect(slide, lx+0.1, 2.02, ws-0.2, 0.03, fill_color=color)
    for bi, b in enumerate(bullets):
        add_text(slide, f"→ {b}", lx+0.14, 2.12+bi*0.85, ws-0.24, 0.8,
                 font_size=11, color=WHITE)

add_text(slide, "Paths can be combined — e.g. PATH A + PATH C for a reaction with an R-group table.",
         0.3, 7.3, 12.7, 0.18, font_size=10, color=LGRAY, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Pipeline Stage 4-5
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
fill_bg(slide)
slide_header(slide, "Pipeline: Stages 4 – 5",
             "Validation → Final structured output")

# Stage 4
add_rect(slide, 0.35, 1.2, 5.9, 5.6, fill_color=MID_BG,
         line_color=ACCENT2, line_w_pt=2)
add_text(slide, "STAGE 4 — Action Observer", 0.5, 1.26, 5.6, 0.45,
         font_size=16, bold=True, color=ACCENT2)
add_rect(slide, 0.5, 1.72, 5.6, 0.03, fill_color=ACCENT2)

obs_items = [
    ("Model",   "GPT-4o-mini / Qwen3-VL"),
    ("Input",   "Image + all tool execution logs"),
    ("Task",    "Validate results; flag errors; signal redo"),
    ("Output",  '{"redo": true/false, "reason": "..."}'),
    ("Disabled","By default in ChemEagle_OS()"),
]
for ii, (lbl, val) in enumerate(obs_items):
    ry = 1.82 + ii * 0.72
    add_text(slide, lbl, 0.55, ry, 1.2, 0.6,
             font_size=11, bold=True, color=ACCENT2)
    add_text(slide, val, 1.8, ry, 4.3, 0.6,
             font_size=11, color=WHITE)

add_rect(slide, 0.4, 4.8, 5.8, 1.0, fill_color=RGBColor(0x3D,0x1A,0x00),
         line_color=WARN, line_w_pt=1)
add_text(slide, "⚠  Architectural gap: orchestrator returns {\"redo\":True} and exits.\n"
                "   The caller must implement the retry loop — the README does not.",
         0.55, 4.85, 5.55, 0.9, font_size=10, color=WARN)

# arrow
add_text(slide, "▶", 6.35, 3.85, 0.55, 0.5,
         font_size=24, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)

# Stage 5
add_rect(slide, 6.95, 1.2, 6.0, 5.6, fill_color=MID_BG,
         line_color=OK, line_w_pt=2)
add_text(slide, "STAGE 5 — Data Structure Agent", 7.1, 1.26, 5.7, 0.45,
         font_size=16, bold=True, color=OK)
add_rect(slide, 7.1, 1.72, 5.7, 0.03, fill_color=OK)

struct_items = [
    ("Model",   "GPT-4o-mini / Qwen3-VL"),
    ("Prompt",  "prompt_final_simple_version.txt"),
    ("Task",    "Merge all tool outputs into one record"),
    ("Converts","Abbreviations → SMILES (from LLM memory)"),
    ("Output",  '{"reactions": [ {reactants, products,\n  conditions, SMILES} ]}'),
]
for ii, (lbl, val) in enumerate(struct_items):
    ry = 1.82 + ii * 0.72
    add_text(slide, lbl, 7.1, ry, 1.35, 0.6,
             font_size=11, bold=True, color=OK)
    add_text(slide, val, 8.5, ry, 4.3, 0.6,
             font_size=11, color=WHITE)

add_rect(slide, 7.0, 4.8, 5.9, 1.0, fill_color=RGBColor(0x0A,0x2A,0x1A),
         line_color=OK, line_w_pt=1)
add_text(slide, "Final output is used directly for downstream cheminformatics,\n"
                "reaction database population, or retrosynthesis tools.",
         7.15, 4.85, 5.65, 0.9, font_size=10, color=OK)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Failure Overview Table
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
fill_bg(slide)
slide_header(slide, "Failure Mode Overview",
             "10 identified failure modes mapped to pipeline stage and primary cause")

failures = [
    ("#",  "Failure",                         "Primary Stage",  "Primary Component",          "Sev"),
    ("1",  "3D / Stereo structures lost",      "Stage 3",        "MolNexTR",                   "HIGH"),
    ("2",  "Complex graphics misclassified",   "Stage 1",        "Planner LLM",                "HIGH"),
    ("3",  "Non-chem PDFs included",           "PDF Extract",    "VisualHeist",                "HIGH"),
    ("4",  "R-group agent mispicked",          "Stage 1",        "Planner LLM",                "MED"),
    ("5",  "Multi-step reactions merged",      "Stage 3",        "RxnIM",                      "MED"),
    ("6",  "Abbreviation SMILES hallucinated", "Stage 3 & 5",    "LLM (no validation)",        "HIGH"),
    ("7",  "Redo signal never executes",        "Stage 4",        "Orchestrator",               "MED"),
    ("8",  "OCR quality degradation",          "Stage 3",        "TesseractOCR / TableParser", "MED"),
    ("9",  "NLP domain mismatch",              "Stage 3",        "ChemRxnExtractor",           "LOW"),
    ("10", "Label absorbed into crop",         "Stage 3",        "MolDetector",                "LOW"),
]

col_x2 = [0.3, 0.78, 4.8, 7.3, 10.8]
col_w2 = [0.44, 3.95, 2.44, 3.44, 1.08]
hdr_y2 = 1.22
row_h2 = 0.46

sev_colors = {"HIGH": DANGER, "MED": WARN, "LOW": OK}

for ri, row in enumerate(failures):
    ry = hdr_y2 + ri * row_h2
    is_header = ri == 0
    bg = ACCENT1 if is_header else (RGBColor(0x1A,0x2C,0x48) if ri%2==0 else MID_BG)
    add_rect(slide, 0.28, ry, 12.75, row_h2-0.03, fill_color=bg)
    for ci, (val, cx, cw) in enumerate(zip(row, col_x2, col_w2)):
        fg = DARK_BG if is_header else WHITE
        bld = is_header or ci == 0
        fsz = 11 if is_header else 10
        if ci == 4 and not is_header:
            sc = sev_colors.get(val, WHITE)
            add_rect(slide, cx+0.04, ry+0.07, cw-0.08, row_h2-0.17,
                     fill_color=RGBColor(int(sc[0]*0.25), int(sc[1]*0.25), int(sc[2]*0.25)),
                     line_color=sc, line_w_pt=1)
            add_text(slide, val, cx+0.04, ry+0.08, cw-0.08, row_h2-0.18,
                     font_size=9, bold=True, color=sc, align=PP_ALIGN.CENTER)
        else:
            add_text(slide, val, cx+0.05, ry+0.06, cw-0.08, row_h2-0.1,
                     font_size=fsz, bold=bld, color=fg)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Top 5 Critical Failures (deep dive)
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
fill_bg(slide)
slide_header(slide, "Top 5 Critical Failures — Deep Dive",
             "The failures with the broadest impact across all pipeline paths")

top5 = [
    ("#1 — 3D/Stereo (MolNexTR)",
     DANGER,
     "MolNexTR is trained on 2D Kekulé structures only. Wedge/dash bonds are "
     "treated as regular bonds; all stereocenters are silently erased. Every "
     "pipeline path routes through MolNexTR, so stereo loss is universal."),
    ("#3 — Non-chem PDF images (VisualHeist)",
     WARN,
     "Florence-2 OD extracts ALL figures with no chemistry filter. Every "
     "photograph, spectrum, and gel in a journal PDF enters the full pipeline, "
     "burning compute and generating garbage output."),
    ("#6 — Abbreviation hallucination (LLM)",
     DANGER,
     "The Data Structure Agent prompt explicitly instructs 'convert to SMILES "
     "from your knowledge'. No RDKit valence check or InChI roundtrip validates "
     "the output — invalid SMILES reach the final JSON silently."),
    ("#2 — Complex graphics (Planner LLM)",
     WARN,
     "Mechanism arrows, spectral baselines, and flow-chart arrows all resemble "
     "reaction arrows. The Planner misclassifies the image and the Plan Observer "
     "(same model, same image) typically agrees."),
    ("#7 — Redo loop never fires (Orchestrator)",
     ACCENT2,
     "When the Action Observer flags a bad result it returns {\"redo\":True}, "
     "but the orchestrator immediately returns this dict to the caller and exits. "
     "The retry loop is the caller's responsibility; the README example omits it."),
]

for i, (title, color, desc) in enumerate(top5):
    row = i // 2 if i < 4 else None
    col = i % 2  if i < 4 else None
    if i < 4:
        lx = 0.35 + col * 6.5
        ty = 1.2  + row * 2.65
        bw, bh = 6.1, 2.55
    else:
        lx, ty, bw, bh = 0.35, 6.5, 12.6, 0.93

    add_rect(slide, lx, ty, bw, bh, fill_color=MID_BG,
             line_color=color, line_w_pt=2)
    add_text(slide, title, lx+0.12, ty+0.08, bw-0.22, 0.35,
             font_size=12, bold=True, color=color)
    add_text(slide, desc, lx+0.12, ty+0.45, bw-0.22, bh-0.52,
             font_size=10, color=LGRAY, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Severity / Frequency chart
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
fill_bg(slide)
slide_header(slide, "Failure Impact Matrix",
             "Severity (x-axis) vs. Frequency (y-axis) — bigger = broader impact")

# grid background
add_rect(slide, 1.4, 1.3, 11.4, 5.75, fill_color=RGBColor(0x0D,0x1E,0x33))

# axes labels
add_text(slide, "HIGH FREQUENCY", 1.4, 1.15, 3.0, 0.25,
         font_size=10, color=LGRAY, bold=True)
add_text(slide, "LOW FREQUENCY", 1.4, 6.9, 3.0, 0.25,
         font_size=10, color=LGRAY, bold=True)
add_text(slide, "LOW SEVERITY", 1.5, 7.1, 2.5, 0.25,
         font_size=10, color=LGRAY, align=PP_ALIGN.LEFT)
add_text(slide, "HIGH SEVERITY", 10.5, 7.1, 2.5, 0.25,
         font_size=10, color=LGRAY, align=PP_ALIGN.RIGHT)

# axis arrows
add_text(slide, "▲", 1.52, 1.28, 0.3, 0.3, font_size=14, color=MGRAY)
add_text(slide, "▶", 12.55, 6.84, 0.3, 0.3, font_size=14, color=MGRAY)

# quadrant labels
quad_style = dict(font_size=9, color=MGRAY, italic=True)
add_text(slide, "Monitor", 1.6, 1.4, 2.0, 0.3, **quad_style)
add_text(slide, "CRITICAL", 10.5, 1.4, 2.2, 0.3, font_size=9, color=DANGER, bold=True, italic=True)
add_text(slide, "Low priority", 1.6, 6.6, 2.5, 0.3, **quad_style)
add_text(slide, "Fix Soon", 10.5, 6.6, 2.0, 0.3, **quad_style)

# dividers
add_rect(slide, 6.9, 1.3, 0.03, 5.75, fill_color=MGRAY)
add_rect(slide, 1.4, 4.18, 11.4, 0.03, fill_color=MGRAY)

# issues: (label, severity 0-1, frequency 0-1, color)
# severity: 0=left(low), 1=right(high)
# frequency: 0=bottom(low), 1=top(high)
issues_plot = [
    ("#1 Stereo",       0.95, 0.92, DANGER),
    ("#6 Abbrev SMILES",0.88, 0.88, DANGER),
    ("#3 Non-chem PDF", 0.75, 0.85, WARN),
    ("#2 Complex gfx",  0.70, 0.65, WARN),
    ("#5 Multi-step",   0.60, 0.60, WARN),
    ("#8 OCR quality",  0.50, 0.62, WARN),
    ("#7 Redo loop",    0.55, 0.30, ACCENT2),
    ("#4 R-group pick", 0.45, 0.35, ACCENT2),
    ("#9 NLP mismatch", 0.25, 0.28, OK),
    ("#10 Label crop",  0.35, 0.20, OK),
]

grid_x0, grid_x1 = 1.4, 12.8
grid_y0, grid_y1 = 7.05, 1.3
gw = grid_x1 - grid_x0
gh = grid_y1 - grid_y0   # negative because y goes up

for label, sev, freq, color in issues_plot:
    px = grid_x0 + sev * gw
    py = grid_y0 + freq * gh   # gh is negative
    # dot
    dot_r = 0.22
    add_rect(slide, px - dot_r, py - dot_r/2, dot_r*2, dot_r,
             fill_color=color, line_color=None)
    add_text(slide, label, px - 0.8, py - 0.38, 1.6, 0.28,
             font_size=8, bold=True, color=color, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Key Findings & Recommendations
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
fill_bg(slide)
slide_header(slide, "Key Findings & Recommendations",
             "Prioritised actions to improve pipeline reliability")

recs = [
    (DANGER, "1", "Replace or retrain MolNexTR with a stereo-aware model",
     "Use RXN-BERT, DECIMER, or Img2Mol trained on 3D/stereo structures. "
     "Add post-processing to re-read wedge/dash bonds from original crop."),
    (DANGER, "2", "Add SMILES validation at every output point",
     "Run RDKit valence check + InChI roundtrip on all generated SMILES. "
     "Reject and flag any string that fails rather than passing it silently."),
    (WARN,   "3", "Insert a chemistry-content classifier before ChemEAGLE",
     "Fine-tune a lightweight ViT or CLIP model to label crops as "
     "'reaction scheme / R-group table / text / other'. Skip non-chemistry images."),
    (WARN,   "4", "Implement the Action Observer retry loop in the orchestrator",
     "Move retry logic inside main.py. Expose max_retries param. "
     "Enable action observer by default in ChemEagle_OS()."),
    (ACCENT2,"5", "Replace TesseractOCR with chemistry-aware OCR",
     "Use DECIMER-Segmentation + MolScribe for structure detection; "
     "use Nougat or GROBID for formula-aware text extraction."),
]

for i, (color, num, title, body) in enumerate(recs):
    ty = 1.25 + i * 1.2
    # number badge
    add_rect(slide, 0.3, ty+0.1, 0.5, 0.75, fill_color=color)
    add_text(slide, num, 0.3, ty+0.1, 0.5, 0.75,
             font_size=20, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    # card
    add_rect(slide, 0.85, ty, 12.1, 1.1, fill_color=MID_BG,
             line_color=color, line_w_pt=1)
    add_text(slide, title, 1.0, ty+0.04, 11.8, 0.38,
             font_size=13, bold=True, color=color)
    add_text(slide, body, 1.0, ty+0.44, 11.8, 0.6,
             font_size=11, color=LGRAY, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════════════════
out_path = "/home/user/ChemEagle/ChemEAGLE_Pipeline_Analysis.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
